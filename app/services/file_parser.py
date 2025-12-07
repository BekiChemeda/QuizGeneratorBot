import io
from typing import List, Tuple
from telebot import TeleBot
from telebot.types import Message
from pymongo.database import Database
from datetime import datetime


MAX_FILE_MB = 20
MAX_FILE_BYTES = MAX_FILE_MB * 1024 * 1024


def _read_pdf_bytes(data: bytes) -> str:
    try:
        import pdfminer.high_level  # type: ignore
        with io.BytesIO(data) as fh:
            return pdfminer.high_level.extract_text(fh) or ""
    except Exception:
        return ""


def _read_docx_bytes(data: bytes) -> str:
    try:
        import docx  # type: ignore
        with io.BytesIO(data) as fh:
            document = docx.Document(fh)
            return "\n".join(p.text for p in document.paragraphs)
    except Exception:
        return ""


def _read_ppt_bytes(data: bytes) -> str:
    try:
        import pptx  # type: ignore
        with io.BytesIO(data) as fh:
            presentation = pptx.Presentation(fh)
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
    except Exception:
        return ""


def _read_txt_bytes(data: bytes) -> str:
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def chunk_text(text: str, max_chars: int = 4000) -> List[str]:
    chunks: List[str] = []
    buf = text.strip()
    while buf:
        chunks.append(buf[:max_chars])
        buf = buf[max_chars:]
    return chunks


def save_file_record(db: Database, user_id: int, file_id: str, file_name: str, file_size: int, mime_type: str) -> str:
    doc = {
        "user_id": user_id,
        "file_id": file_id,
        "file_name": file_name,
        "file_size": file_size,
        "mime_type": mime_type,
        "created_at": datetime.now(),
    }
    res = db["files"].insert_one(doc)
    return str(res.inserted_id)


def fetch_and_parse_file(bot: TeleBot, db: Database, message: Message) -> Tuple[str, str]:
    """Return (content_text, file_name). Raises ValueError on invalid/too-big or unknown type."""
    if not message.document:
        raise ValueError("No document attached")
    doc = message.document
    if doc.file_size and doc.file_size > MAX_FILE_BYTES:
        raise ValueError("File exceeds 20 MB. Please split it and send again.")
    file_info = bot.get_file(doc.file_id)
    file_bytes = bot.download_file(file_info.file_path)
    if doc.mime_type and "pdf" in doc.mime_type.lower():
        text = _read_pdf_bytes(file_bytes)
    elif doc.file_name and doc.file_name.lower().endswith(".pdf"):
        text = _read_pdf_bytes(file_bytes)
    elif doc.mime_type and ("word" in doc.mime_type.lower() or "docx" in doc.mime_type.lower()):
        text = _read_docx_bytes(file_bytes)
    elif doc.file_name and doc.file_name.lower().endswith(".docx"):
        text = _read_docx_bytes(file_bytes)
    elif doc.mime_type and ("powerpoint" in doc.mime_type.lower() or "presentation" in doc.mime_type.lower()) or (doc.file_name and doc.file_name.lower().endswith((".ppt", ".pptx"))):
        text = _read_ppt_bytes(file_bytes)
    else:
        # treat as text
        text = _read_txt_bytes(file_bytes)
    if not text.strip():
        raise ValueError("Failed to parse file content.")
    save_file_record(db, message.from_user.id, doc.file_id, doc.file_name or "file", doc.file_size or 0, doc.mime_type or "")
    return text, doc.file_name or "file"

